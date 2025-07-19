import os
import json
import asyncio
import re
import random
from fastapi import FastAPI, HTTPException, File, UploadFile, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import List
import google.generativeai as genai
from dotenv import load_dotenv
from supabase import create_client, Client
import uuid
import io
from pptx import Presentation
from pptx.util import Inches

# --- Environment Variable Setup ---
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY") # Use service role key

# --- FastAPI App Initialization ---
app = FastAPI(
    title="Pitch Deck PPTX Generator API",
    description="An API to generate and save pitch decks as .pptx files.",
    version="9.0.0"
)

# --- CORS Configuration ---
origins = ["*"] # Or specify your front-end domains
app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- API Client Initialization ---
genai.configure(api_key=GEMINI_API_KEY)
supabase: Client = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)

# --- Pydantic Models & Helper Functions ---
class PitchDeckSlide(BaseModel):
    title: str
    content: str

class PitchDeckContent(BaseModel):
    company_name: str
    tagline: str
    slides: List[PitchDeckSlide]

def create_pitch_deck_prompt(idea_text: str) -> str:
    return f"""
    You are an expert startup consultant. Based on the user's idea: "{idea_text}", generate the content for a 5-slide pitch deck.
    Return a single, valid JSON object with the structure: {{"company_name": "...", "tagline": "...", "slides": [{{"title": "...", "content": "..."}}]}}.
    The 5 slides MUST be: 1. The Problem, 2. Our Solution, 3. Business Model, 4. Traction & Milestones, 5. The Ask.
    """

async def _generate_deck_content(idea_text: str) -> dict:
    model = genai.GenerativeModel('gemini-2.5-flash')
    prompt = create_pitch_deck_prompt(idea_text)
    try:
        response = await model.generate_content_async(prompt)
        match = re.search(r'\{.*\}', response.text, re.DOTALL)
        if not match:
            raise ValueError("No valid JSON object found in the AI's response.")
        json_string = match.group(0)
        return json.loads(json_string)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI content generation failed: {e}")

def _create_pitch_deck_pptx(deck_content: dict) -> io.BytesIO:
    """Creates a new PowerPoint presentation from deck content."""
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = deck_content.get("company_name", "Company Name")
    subtitle.text = deck_content.get("tagline", "Tagline")

    # Content Slides
    for slide_data in deck_content.get("slides", []):
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = slide_data.get("title", "Slide Title")
        body.text = slide_data.get("content", "Slide content...")

    # Save to buffer
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

async def _process_and_save_idea(final_idea_text: str):
    try:
        print("Step 2: Generating deck text content...")
        deck_content = await _generate_deck_content(final_idea_text)
        print("Step 2 Complete.")

        print("Step 4: Creating PPTX document...")
        pptx_buffer = _create_pitch_deck_pptx(deck_content)
        pptx_bytes = pptx_buffer.getvalue()
        print("Step 4 Complete.")
        
        print("Step 5: Uploading PPTX to Supabase Storage...")
        file_path = f"public/{uuid.uuid4()}.pptx"
        file_options = {"content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}
        supabase.storage.from_("pitch-decks").upload(file=pptx_bytes, path=file_path, file_options=file_options)
        file_url = supabase.storage.from_("pitch-decks").get_public_url(file_path)
        print("Step 5 Complete.")

        print("Step 6: Saving metadata to database...")
        deck_to_save = {
            "company_name": deck_content.get("company_name"),
            "tagline": deck_content.get("tagline"),
            "original_idea": final_idea_text,
            "pdf_url": file_url, # Consider renaming this field in your database table
            "storage_path": file_path
        }
        db_response = supabase.table("pitch_decks").insert(deck_to_save).execute()
        print("Step 6 Complete.")
        
        return db_response.data[0]
    except Exception as e:
        print(f"\n--- ERROR DURING DECK GENERATION ---")
        print(f"An unexpected error occurred: {e}")
        print("--- END OF ERROR ---")
        raise HTTPException(status_code=500, detail=str(e))

# --- API Endpoints ---
@app.post("/decks/text", tags=["Pitch Decks"])
async def generate_from_text(idea_text: str = Form(...)):
    if len(idea_text) < 20:
        raise HTTPException(status_code=400, detail="Text idea is too short.")
    new_deck_data = await _process_and_save_idea(idea_text)
    return JSONResponse(status_code=201, content=new_deck_data)

@app.post("/decks/audio", tags=["Pitch Decks"])
async def generate_from_audio(audio_file: UploadFile = File(...)):
    transcription_model = genai.GenerativeModel('gemini-1.5-flash-latest')
    audio_bytes = await audio_file.read()
    audio_part = {"mime_type": audio_file.content_type, "data": audio_bytes}
    response = await transcription_model.generate_content_async(["Transcribe this startup idea.", audio_part])
    final_idea_text = response.text.strip()
    new_deck_data = await _process_and_save_idea(final_idea_text)
    return JSONResponse(status_code=201, content=new_deck_data)

@app.get("/decks", tags=["Pitch Decks"])
async def get_all_decks():
    response = supabase.table("pitch_decks").select("*").order("created_at", desc=True).execute()
    return response.data

@app.delete("/decks/{deck_id}", tags=["Pitch Decks"])
async def delete_deck(deck_id: int):
    select_res = supabase.table("pitch_decks").select("storage_path").eq("id", deck_id).single().execute()
    if not select_res.data:
        raise HTTPException(status_code=404, detail="Deck not found.")
    storage_path = select_res.data['storage_path']
    supabase.storage.from_("pitch-decks").remove([storage_path])
    supabase.table("pitch_decks").delete().eq("id", deck_id).execute()
    return {"status": "success", "message": f"Deck {deck_id} deleted."}